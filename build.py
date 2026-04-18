#!/usr/bin/env python3
"""Extract the SOM PowerPoint deck into structured JSON + image assets for the
static viewer under `docs/`.

Slides are 960x540 (logical px). Every shape is emitted with absolute
coordinates in that coordinate system so the viewer can scale the slide to
any viewport size with a single CSS transform.
"""

from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parent
PPTX = ROOT / "SOM_PC_26_Slides.pptx"
DOCS = ROOT / "docs"
MEDIA = DOCS / "media"

SLIDE_W = 960
SLIDE_H = 540
EMU_PER_PX = 9525  # at 96 DPI

# Placeholder typography defaults (px). The deck doesn't set explicit run-level
# font sizes on most runs, so we infer reasonable defaults from the placeholder
# type and let explicit overrides win.
PH_DEFAULTS = {
    "ctrTitle": {"size": 40, "bold": True, "align": "center", "color": "#1F2937"},
    "title":    {"size": 32, "bold": True, "align": "left",   "color": "#1F2937"},
    "subTitle": {"size": 20, "bold": False, "align": "center", "color": "#4B5563"},
    "body":     {"size": 18, "bold": False, "align": "left",   "color": "#1F2937"},
    "obj":      {"size": 18, "bold": False, "align": "left",   "color": "#1F2937"},
}
DEFAULT_TEXT = {"size": 16, "bold": False, "align": "left", "color": "#1F2937"}


def emu_to_px(emu: int | None) -> float:
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_PX, 2)


def align_to_css(a) -> str:
    if a is None:
        return "left"
    s = str(a)
    if "CENTER" in s:
        return "center"
    if "RIGHT" in s:
        return "right"
    if "JUSTIFY" in s:
        return "justify"
    return "left"


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    try:
        return "#" + str(rgb).upper()
    except Exception:
        return None


def placeholder_type(shape) -> str | None:
    try:
        if shape.is_placeholder:
            t = shape.placeholder_format.type
            if t is not None:
                return str(t).split(".")[-1].split(" ")[0].lower()
    except Exception:
        pass
    return None


def collect_text(shape) -> list[dict]:
    """Return a list of paragraph dicts: {align, level, runs:[{text, size, bold, italic, color, underline}]}"""
    paragraphs = []
    tf = shape.text_frame
    for para in tf.paragraphs:
        runs = []
        for run in para.runs:
            font = run.font
            color = None
            try:
                color = rgb_to_hex(font.color.rgb)
            except Exception:
                color = None
            size = None
            if font.size is not None:
                # pptx Pt -> px ~ *1.333
                size = round(font.size.pt * 1.333, 1)
            # PowerPoint uses \v (vertical tab, 0x0b) as a soft line break and
            # \t as a tab stop. Normalize to HTML-friendly forms.
            text = (run.text or "").replace("\v", "\n").replace("\t", "\u00a0\u00a0\u00a0\u00a0")
            runs.append({
                "text": text,
                "size": size,
                "bold": font.bold,
                "italic": font.italic,
                "underline": bool(font.underline) if font.underline is not None else None,
                "color": color,
                "name": font.name,
            })
        # Paragraphs with no runs still contribute a blank line
        if not runs and para.text == "":
            runs = [{"text": ""}]
        paragraphs.append({
            "align": align_to_css(para.alignment),
            "level": para.level or 0,
            "runs": runs,
        })
    return paragraphs


def slide_bg(slide) -> str | None:
    # Try to read solid fill on the slide background
    try:
        bg = slide.background.fill
        if bg.type is not None and str(bg.type).endswith("SOLID"):
            return rgb_to_hex(bg.fore_color.rgb)
    except Exception:
        pass
    return None


def extract_shapes(shapes, out_list: list[dict]):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # groups: recurse; positions inside groups in pptx are in the group's
            # coordinate space but python-pptx's shape.left/top on grouped shapes
            # already reflects the absolute position after transform, so we just
            # recurse through the group's shapes.
            extract_shapes(shape.shapes, out_list)
            continue

        entry = {
            "x": emu_to_px(shape.left),
            "y": emu_to_px(shape.top),
            "w": emu_to_px(shape.width),
            "h": emu_to_px(shape.height),
            "name": shape.name,
        }

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = img.ext or "png"
            sha = img.sha1
            filename = f"img_{sha}.{ext}"
            target = MEDIA / filename
            if not target.exists():
                target.write_bytes(img.blob)
            entry["type"] = "image"
            entry["src"] = f"media/{filename}"
            out_list.append(entry)
            continue

        if shape.has_text_frame:
            paras = collect_text(shape)
            has_text = any(r.get("text") for p in paras for r in p.get("runs", []))
            if has_text:
                ph = placeholder_type(shape)
                defaults = PH_DEFAULTS.get(ph, DEFAULT_TEXT)
                entry["type"] = "text"
                entry["paragraphs"] = paras
                entry["defaults"] = defaults
                # Container padding (pt-ish). Slide master usually has ~8pt inset.
                entry["pad"] = 8
                out_list.append(entry)
                continue

        # Non-text auto shapes with a solid fill — render a colored rectangle
        try:
            fill = shape.fill
            if fill.type is not None and str(fill.type).endswith("SOLID"):
                c = rgb_to_hex(fill.fore_color.rgb)
                if c:
                    entry["type"] = "rect"
                    entry["color"] = c
                    out_list.append(entry)
        except Exception:
            pass


def build() -> None:
    DOCS.mkdir(exist_ok=True)
    MEDIA.mkdir(exist_ok=True)

    prs = Presentation(PPTX)
    slides_out = []
    for idx, slide in enumerate(prs.slides, start=1):
        elements: list[dict] = []
        extract_shapes(slide.shapes, elements)
        slides_out.append({
            "index": idx,
            "layout": slide.slide_layout.name,
            "bg": slide_bg(slide),
            "elements": elements,
        })

    payload = {
        "width": SLIDE_W,
        "height": SLIDE_H,
        "slides": slides_out,
        "title": "SOM Precourse 2026",
    }
    (DOCS / "slides.json").write_text(json.dumps(payload, ensure_ascii=False, indent=1))
    print(f"wrote {len(slides_out)} slides + {len(list(MEDIA.iterdir()))} images")


if __name__ == "__main__":
    build()
